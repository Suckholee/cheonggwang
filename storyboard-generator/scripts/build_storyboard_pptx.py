#!/usr/bin/env python3
"""
build_storyboard_pptx.py — Storyboard JSON → 통합 .pptx

슬라이드 순서:
  1. 표지
  2. 개정 이력
  3. IA 흐름도 (도형으로 직접 렌더; 카테고리 색상 + 박스 항목 + 화살표 + 범례)
  4 ~ N. 화면별 상세 (5분할 레이아웃)

사용법:
    python build_storyboard_pptx.py --input storyboard.json --output 스토리보드.pptx
    python build_storyboard_pptx.py --input storyboard.json --output 스토리보드.pptx --ia-render image
                                                                        (default: shapes)

의존성: python-pptx, Pillow (이미지 fit용)
    pip install python-pptx Pillow
"""
from __future__ import annotations

import argparse
import json
import os
import sys
from typing import Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError as e:
    print("ERROR: python-pptx와 Pillow가 필요합니다. `pip install python-pptx Pillow`", file=sys.stderr)
    raise

try:
    from PIL import Image
except ImportError:
    Image = None  # 일부 기능만 동작


# ---------------------------------------------------------------------------
# 토큰
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

CLR_TEXT = RGBColor(0x11, 0x18, 0x27)
CLR_TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
CLR_BORDER = RGBColor(0x94, 0xA3, 0xB8)
CLR_HEADER_GRAY = RGBColor(0x6B, 0x72, 0x80)
CLR_TABLE_HEADER = RGBColor(0x9C, 0xA3, 0xAF)
CLR_BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
CLR_RED = RGBColor(0xDC, 0x26, 0x26)
CLR_BLUE = RGBColor(0x3B, 0x82, 0xF6)
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_GRAY_DARK = RGBColor(0x1F, 0x29, 0x37)

DEFAULT_PALETTE = {
    "auth":    "#F59E0B",
    "main":    "#FACC15",
    "extra":   "#9CA3AF",
    "config":  "#94A3B8",
    "admin":   "#DC2626",
    "billing": "#2563EB",
    "content": "#14B8A6",
}


def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# 도형/텍스트 헬퍼
# ---------------------------------------------------------------------------
def add_text(slide, x, y, w, h, text, *, size=12, bold=False, color=CLR_TEXT, align=None, anchor=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    if anchor:
        tf.vertical_anchor = anchor
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_multiline(slide, x, y, w, h, lines: list[tuple[str, dict]], *, size_default=10, color_default=CLR_TEXT):
    """여러 줄을 하나의 textbox에 넣는다. lines: [(text, {size?, bold?, color?, indent?}) ...]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    first = True
    for text, opts in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if opts.get("indent"):
            p.level = int(opts["indent"])
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", size_default))
        run.font.bold = bool(opts.get("bold", False))
        run.font.color.rgb = opts.get("color", color_default)
    return tb


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=Pt(0.75), corner=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_w
    shape.shadow.inherit = False
    if corner:
        # 라운드 코너 비율 조정: ROUNDED_RECTANGLE은 자동, 별도 조정 없이 그대로 사용
        pass
    # 텍스트 비활성
    shape.text_frame.text = ""
    return shape


def add_oval(slide, x, y, w, h, *, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def add_arrow(slide, x1, y1, x2, y2, *, color=RGBColor(0x6B, 0x72, 0x80), width=Pt(1)):
    """단순 직선 화살표 (CONNECTOR_STRAIGHT). PPTX 도형으로."""
    from pptx.util import Emu
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = width
    # 화살촉 추가 (XML 직접 조작)
    spPr = line.line._get_or_add_ln()
    tail = etree.SubElement(spPr, qn("a:tailEnd"))
    tail.set("type", "triangle")
    return line


# ---------------------------------------------------------------------------
# 슬라이드 1 — 표지
# ---------------------------------------------------------------------------
def slide_cover(prs: Presentation, meta: dict) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    title = meta.get("service") or meta.get("title", "Storyboard")
    subtitle = "스토리보드 (화면 정의서)"

    add_text(s, Inches(0), Inches(2.4), SLIDE_W, Inches(0.9), title,
             size=44, bold=True, color=CLR_TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(3.4), SLIDE_W, Inches(0.6), subtitle,
             size=22, color=CLR_TEXT_MUTED, align=PP_ALIGN.CENTER)

    # 구분선
    add_rect(s, Inches(5), Inches(4.2), Inches(3.333), Pt(2), fill=CLR_BORDER, line=None)

    info_lines = []
    info_lines.append(("v" + str(meta.get("version", "0.1")) + " · " + str(meta.get("created_at", "")), {"size": 14, "color": CLR_TEXT}))
    if meta.get("author"):
        info_lines.append(("작성: " + str(meta["author"]), {"size": 12, "color": CLR_TEXT_MUTED}))
    if meta.get("client"):
        info_lines.append(("클라이언트: " + str(meta["client"]), {"size": 12, "color": CLR_TEXT_MUTED}))
    platform_kr = {
        "mobile_app": "모바일 앱",
        "responsive_web": "반응형 웹",
        "desktop_web": "데스크톱 웹",
        "multi": "모바일 + 웹",
    }.get(meta.get("platform"), str(meta.get("platform", "")))
    if platform_kr:
        info_lines.append(("플랫폼: " + platform_kr, {"size": 12, "color": CLR_TEXT_MUTED}))
    if meta.get("purpose"):
        info_lines.append(("용도: " + str(meta["purpose"]), {"size": 12, "color": CLR_TEXT_MUTED}))

    tb = slide_paragraphs_centered(s, Inches(0), Inches(4.5), SLIDE_W, Inches(2.0), info_lines)


def slide_paragraphs_centered(slide, x, y, w, h, lines):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for text, opts in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", 12))
        run.font.bold = bool(opts.get("bold", False))
        if opts.get("color"):
            run.font.color.rgb = opts["color"]
    return tb


# ---------------------------------------------------------------------------
# 슬라이드 2 — 개정 이력
# ---------------------------------------------------------------------------
def slide_revision_history(prs: Presentation, history: list[dict]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(s, Inches(0), Inches(0.4), SLIDE_W, Inches(0.6),
             "개정 이력 (Revision History)", size=24, bold=True, align=PP_ALIGN.CENTER)

    if not history:
        add_text(s, Inches(0), Inches(3.5), SLIDE_W, Inches(0.5),
                 "(기록 없음)", size=14, color=CLR_TEXT_MUTED, align=PP_ALIGN.CENTER)
        return

    # 표
    rows = len(history) + 1
    cols = 4
    left = Inches(1.5)
    top = Inches(1.4)
    width = Inches(10.333)
    height = Inches(min(5.5, 0.45 + 0.4 * len(history)))

    table_shape = s.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # 컬럼 너비
    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(1.6)
    table.columns[2].width = Inches(1.8)
    table.columns[3].width = Inches(5.733)

    headers = ["Version", "Date", "Author", "Summary"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(11)
            run.font.color.rgb = CLR_WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = CLR_HEADER_GRAY

    for r, rev in enumerate(history, start=1):
        vals = [
            str(rev.get("version", "")),
            str(rev.get("date", "")),
            str(rev.get("author", "")),
            str(rev.get("summary", "")),
        ]
        for c, v in enumerate(vals):
            cell = table.cell(r, c)
            cell.text = v
            para = cell.text_frame.paragraphs[0]
            for run in para.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = CLR_TEXT


# ---------------------------------------------------------------------------
# 슬라이드 3 — IA 흐름도 (도형으로 직접)
# ---------------------------------------------------------------------------
def slide_ia(prs: Presentation, data: dict, mode: str = "shapes", ia_image_path: str | None = None) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    ia = data.get("ia", {})

    # 제목
    add_text(s, Inches(0.4), Inches(0.3), Inches(8), Inches(0.4),
             "IA / Site Map", size=20, bold=True)
    add_text(s, Inches(0.4), Inches(0.7), Inches(8), Inches(0.3),
             data.get("meta", {}).get("title", ""), size=10, color=CLR_TEXT_MUTED)

    # 범례
    categories = ia.get("categories", [])
    if categories:
        legend_x = Inches(8.5)
        legend_y = Inches(0.3)
        legend_w = Inches(4.5)
        rows = (len(categories) + 1) // 2
        legend_h = Inches(0.3 + 0.3 * rows)
        add_rect(s, legend_x, legend_y, legend_w, legend_h, fill=CLR_WHITE, line=CLR_BORDER, corner=True)
        for i, cat in enumerate(categories):
            cx = legend_x + Inches(0.15 + (i % 2) * 2.2)
            cy = legend_y + Inches(0.1 + (i // 2) * 0.3)
            color_hex = cat.get("color", DEFAULT_PALETTE.get(cat.get("id"), "#9CA3AF"))
            add_rect(s, cx, cy, Inches(0.18), Inches(0.18), fill=hex_to_rgb(color_hex), line=None)
            add_text(s, cx + Inches(0.25), cy - Inches(0.02), Inches(2), Inches(0.25),
                     cat.get("name", cat.get("id", "")), size=10)

    if mode == "image" and ia_image_path and os.path.exists(ia_image_path):
        # 이미지 모드: SVG → PNG 변환된 IA 이미지를 박는다
        s.shapes.add_picture(ia_image_path, Inches(0.4), Inches(1.2),
                             width=Inches(12.5), height=Inches(6.0))
        return

    # shapes 모드: 박스를 직접 그린다
    boxes = ia.get("boxes", [])
    connections = ia.get("connections", [])
    if not boxes:
        add_text(s, Inches(0.4), Inches(3), SLIDE_W, Inches(0.5),
                 "(IA boxes 비어 있음)", size=14, color=CLR_TEXT_MUTED, align=PP_ALIGN.CENTER)
        return

    cat_color: dict[str, str] = {}
    for c in categories:
        cat_color[c.get("id")] = c.get("color", DEFAULT_PALETTE.get(c.get("id"), "#9CA3AF"))
    for b in boxes:
        cid = b.get("category")
        if cid and cid not in cat_color:
            cat_color[cid] = DEFAULT_PALETTE.get(cid, "#9CA3AF")

    # 박스 측정 (간단 버전: 항목 수 × 18px + 헤더)
    def measure(b):
        items = b.get("items", []) or []
        h = 0
        for it in items:
            h += 18
            for _ in (it.get("subitems") or []):
                h += 14
        oqs = b.get("open_questions", []) or []
        for q in oqs:
            chars = len(str(q))
            lines = max(1, (chars // 28))
            h += 12 * (lines + 1) + 4
        # PPT inches로 변환 (대략 px/96 → inches)
        return max(Inches(1.0), Inches(0.35 + h / 72))

    # 자동 배치 (간단 버전)
    box_h: dict[str, Any] = {b["id"]: measure(b) for b in boxes}
    # 카테고리 순서대로 col 그룹화 + 각 카테고리는 박스 등장 순서대로 row
    placement: dict[str, tuple[int, int]] = {}
    incoming: dict[str, int] = {b["id"]: 0 for b in boxes}
    for c in connections:
        if c.get("to") in incoming:
            incoming[c["to"]] += 1

    # depth 계산 (단순 BFS)
    box_ids = [b["id"] for b in boxes]
    depth = {bid: 0 for bid in box_ids}
    box_dict = {b["id"]: b for b in boxes}
    # 진입 = incoming 0인 박스
    queue = [bid for bid in box_ids if incoming[bid] == 0]
    visited = set(queue)
    while queue:
        cur = queue.pop(0)
        for c in connections:
            if c.get("from") == cur and c.get("to") in box_dict:
                t = c["to"]
                depth[t] = max(depth[t], depth[cur] + 1)
                if t not in visited:
                    visited.add(t)
                    queue.append(t)

    # 사용자가 row/col을 명시했으면 우선시
    for b in boxes:
        if "col" in b and "row" in b:
            placement[b["id"]] = (int(b["col"]), int(b["row"]))

    # 미배치 박스 자동 배치
    used: set[tuple[int, int]] = set(placement.values())
    for b in boxes:
        if b["id"] in placement:
            continue
        col = depth.get(b["id"], 0)
        row = 0
        while (col, row) in used:
            row += 1
        placement[b["id"]] = (col, row)
        used.add((col, row))

    # 좌표 환산
    diagram_left = Inches(0.4)
    diagram_top = Inches(1.3)
    box_w = Inches(2.0)
    col_gap = Inches(0.4)
    row_gap = Inches(0.25)

    # col별 x
    cols = sorted({c for c, _ in placement.values()})
    col_x = {c: diagram_left + Inches(0) + (box_w + col_gap) * i for i, c in enumerate(cols)}

    # 같은 col 안에서 row순으로 y 누적
    rows_by_col: dict[int, list[tuple[int, str]]] = {}
    for bid, (c, r) in placement.items():
        rows_by_col.setdefault(c, []).append((r, bid))
    box_pos: dict[str, tuple[Any, Any, Any, Any]] = {}  # (x, y, w, h)
    for c, lst in rows_by_col.items():
        lst.sort(key=lambda t: t[0])
        cur_y = diagram_top
        for _, bid in lst:
            h = box_h[bid]
            box_pos[bid] = (col_x[c], cur_y, box_w, h)
            cur_y = cur_y + h + row_gap

    # 박스 그리기
    for bid, (x, y, w, h) in box_pos.items():
        b = box_dict[bid]
        cat_id = b.get("category", "")
        header_color_hex = cat_color.get(cat_id, "#9CA3AF")
        header_color = hex_to_rgb(header_color_hex)

        # 본체
        line_color = CLR_BORDER
        body = add_rect(s, x, y, w, h, fill=CLR_WHITE, line=line_color, corner=True)
        if b.get("auto_filled"):
            # 점선은 python-pptx에서 직접 지원이 까다로움; 라인 색만 살짝 다르게
            body.line.color.rgb = RGBColor(0xC0, 0xC0, 0xC0)

        # 헤더
        header_h = Inches(0.3)
        add_rect(s, x, y, w, header_h, fill=header_color, line=None)
        add_text(s, x + Inches(0.1), y + Inches(0.04), w - Inches(0.2), header_h,
                 b.get("title", b.get("id", "")), size=11, bold=True, color=CLR_WHITE)

        # 항목
        cy = y + header_h + Inches(0.05)
        for it in (b.get("items") or []):
            no = str(it.get("no", ""))
            label = str(it.get("label", ""))
            add_text(s, x + Inches(0.08), cy, Inches(0.25), Inches(0.22), no, size=9, bold=True)
            add_text(s, x + Inches(0.32), cy, w - Inches(0.4), Inches(0.22), label, size=9)
            cy += Inches(0.22)
            for sub in (it.get("subitems") or []):
                add_text(s, x + Inches(0.4), cy, w - Inches(0.5), Inches(0.18),
                         "• " + str(sub), size=8, color=CLR_TEXT_MUTED)
                cy += Inches(0.18)
        # open_questions
        if b.get("open_questions"):
            cy += Inches(0.05)
            for q in b["open_questions"]:
                add_text(s, x + Inches(0.1), cy, w - Inches(0.2), Inches(0.5),
                         "※ " + str(q), size=8, color=CLR_RED)
                cy += Inches(0.2)

    # 화살표
    for c in connections:
        f = c.get("from")
        t = c.get("to")
        if f not in box_pos or t not in box_pos:
            continue
        fx, fy, fw, fh = box_pos[f]
        tx, ty, tw, th = box_pos[t]

        style = c.get("style", "secondary")
        color = CLR_GRAY_DARK if style == "primary" else CLR_TEXT_MUTED
        width = Pt(1.5) if style == "primary" else Pt(0.75)

        # 좌→우 케이스
        if fx + fw <= tx:
            x1 = fx + fw
            y1 = fy + fh / 2
            x2 = tx
            y2 = ty + th / 2
        elif fx >= tx + tw:
            x1 = fx
            y1 = fy + fh / 2
            x2 = tx + tw
            y2 = ty + th / 2
        elif fy + fh <= ty:
            x1 = fx + fw / 2
            y1 = fy + fh
            x2 = tx + tw / 2
            y2 = ty
        else:
            x1 = fx + fw / 2
            y1 = fy
            x2 = tx + tw / 2
            y2 = ty + th
        add_arrow(s, x1, y1, x2, y2, color=color, width=width)

        if c.get("label"):
            mx = (x1 + x2) / 2
            my = (y1 + y2) / 2
            add_text(s, mx - Inches(0.5), my - Inches(0.12), Inches(1.0), Inches(0.24),
                     str(c["label"]), size=8, color=CLR_TEXT, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# 슬라이드 4~N — 화면별 상세
# ---------------------------------------------------------------------------
def slide_screen_detail(prs: Presentation, screen: dict, doc_root: str, presets: dict) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더: 섹션 번호 + 제목
    section_text = ""
    if screen.get("section_no"):
        section_text = f"{screen['section_no']}. "
    section_text += screen.get("section_title", "와이어프레임 & 기능 디스크립션")

    add_text(s, Inches(0.4), Inches(0.25), Inches(10), Inches(0.45),
             section_text, size=20, bold=True)

    # PATH 박스
    path_y = Inches(0.85)
    path_left = Inches(0.4)
    path_label_w = Inches(0.8)
    path_value_w = Inches(7.0)
    path_h = Inches(0.35)
    add_rect(s, path_left, path_y, path_label_w, path_h, fill=CLR_HEADER_GRAY, line=CLR_BORDER)
    add_text(s, path_left, path_y + Inches(0.04), path_label_w, path_h,
             "PATH", size=10, bold=True, color=CLR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, path_left + path_label_w, path_y, path_value_w, path_h, fill=CLR_WHITE, line=CLR_BORDER)
    add_text(s, path_left + path_label_w + Inches(0.1), path_y + Inches(0.05), path_value_w - Inches(0.2), path_h,
             screen.get("path", screen.get("title", "")), size=10)

    # 좌측: 컴포넌트 상태 변형 stack
    left_x = Inches(0.4)
    left_y = Inches(1.4)
    left_w = Inches(2.4)
    left_h = Inches(5.8)
    states = screen.get("component_states", [])
    if states:
        each_h = min(Inches(0.85), left_h / max(len(states), 1))
        cur_y = left_y
        for st in states:
            label = st.get("label", st.get("preset_id", ""))
            preset_id = st.get("preset_id")
            preset = presets.get(preset_id, {}) if preset_id else st.get("inline", {})
            add_text(s, left_x, cur_y, left_w, Inches(0.22),
                     label, size=10, bold=True, color=CLR_BLUE)
            mini_y = cur_y + Inches(0.25)
            mini_h = Inches(0.45)
            render_state_mini(s, left_x, mini_y, left_w, mini_h, preset)
            cur_y += each_h

    # 중앙: 메인 목업
    mid_x = Inches(2.95)
    mid_y = Inches(1.4)
    mid_w = Inches(5.4)
    mid_h = Inches(5.8)
    add_rect(s, mid_x, mid_y, mid_w, mid_h, fill=CLR_BG_LIGHT, line=CLR_BORDER, corner=True)

    main_mockup = screen.get("main_mockup", {})
    img_actual_box = (mid_x, mid_y, mid_w, mid_h)
    if main_mockup.get("kind") == "image" and main_mockup.get("path"):
        img_path = main_mockup["path"]
        if not os.path.isabs(img_path):
            img_path = os.path.join(doc_root, img_path)
        if os.path.exists(img_path):
            try:
                if Image is not None:
                    img = Image.open(img_path)
                    iw, ih = img.size
                    box_ratio = mid_w / mid_h
                    img_ratio = iw / ih
                    if img_ratio > box_ratio:
                        new_w = mid_w - Inches(0.2)
                        new_h = new_w / img_ratio
                    else:
                        new_h = mid_h - Inches(0.2)
                        new_w = new_h * img_ratio
                    px = mid_x + (mid_w - new_w) / 2
                    py = mid_y + (mid_h - new_h) / 2
                    s.shapes.add_picture(img_path, px, py, width=new_w, height=new_h)
                    img_actual_box = (px, py, new_w, new_h)
                else:
                    s.shapes.add_picture(img_path, mid_x + Inches(0.1), mid_y + Inches(0.1),
                                         width=mid_w - Inches(0.2), height=mid_h - Inches(0.2))
                    img_actual_box = (mid_x + Inches(0.1), mid_y + Inches(0.1), mid_w - Inches(0.2), mid_h - Inches(0.2))
            except Exception as e:
                add_text(s, mid_x + Inches(0.5), mid_y + Inches(2.5), mid_w - Inches(1), Inches(0.5),
                         f"(이미지 로드 실패: {e})", size=10, color=CLR_TEXT_MUTED, align=PP_ALIGN.CENTER)
        else:
            add_text(s, mid_x + Inches(0.5), mid_y + Inches(2.5), mid_w - Inches(1), Inches(0.5),
                     f"(이미지 없음: {img_path})", size=10, color=CLR_TEXT_MUTED, align=PP_ALIGN.CENTER)
    else:
        add_text(s, mid_x + Inches(0.5), mid_y + Inches(2.5), mid_w - Inches(1), Inches(0.5),
                 "(메인 목업: 디자인 시안 업로드 또는 wireframe 자동 생성 필요)",
                 size=10, color=CLR_TEXT_MUTED, align=PP_ALIGN.CENTER)

    # 어노테이션 마커 오버레이
    img_x, img_y, img_w, img_h = img_actual_box
    for ann in main_mockup.get("annotations", []) or []:
        nx = float(ann.get("x", 0.5))
        ny = float(ann.get("y", 0.5))
        marker = str(ann.get("marker", "?"))
        kind = ann.get("kind", "number")
        marker_size = Inches(0.32)
        cx = img_x + img_w * nx - marker_size / 2
        cy = img_y + img_h * ny - marker_size / 2
        fill_color = CLR_RED if kind == "number" else CLR_BLUE
        add_oval(s, cx, cy, marker_size, marker_size, fill=fill_color, line=None)
        add_text(s, cx, cy + Inches(0.02), marker_size, marker_size,
                 marker, size=11, bold=True, color=CLR_WHITE, align=PP_ALIGN.CENTER)

    # 우측: Function / Button / Validation
    right_x = Inches(8.5)
    right_w = Inches(4.6)
    right_y = Inches(1.4)

    # Function Description
    fns = screen.get("function_descriptions", [])
    if fns:
        right_y = render_desc_table(s, right_x, right_y, right_w, "Function Description", fns)
        right_y += Inches(0.15)

    # Button Description
    btns = screen.get("button_descriptions", [])
    if btns:
        right_y = render_desc_table(s, right_x, right_y, right_w, "Button Description", btns)
        right_y += Inches(0.15)

    # Validation 표
    vals = screen.get("validation_table", [])
    if vals:
        render_validation_table(s, right_x, right_y, right_w, vals)


def render_state_mini(s, x, y, w, h, preset: dict):
    """좌측 상태 변형 미니어처를 단순화해 그린다."""
    kind = preset.get("kind", "input")
    if kind == "input":
        state = preset.get("state", "basic")
        # 기본 입력 박스
        add_rect(s, x, y + h - Inches(0.1), w, Pt(1), fill=CLR_BORDER, line=None)
        # 상태별 라벨/장식
        label = preset.get("placeholder", "")
        color = CLR_TEXT_MUTED
        if state == "focused":
            add_rect(s, x, y + h - Inches(0.06), w, Pt(2), fill=CLR_BLUE, line=None)
            label = preset.get("placeholder", "비밀번호") + " |"
        elif state == "typing":
            label = preset.get("text", "germ")
            color = CLR_TEXT
            # 우측 X
            add_oval(s, x + w - Inches(0.18), y + h - Inches(0.28), Inches(0.14), Inches(0.14),
                     fill=RGBColor(0xCB, 0xD5, 0xE1), line=None)
        elif state == "masked":
            label = "●●●m"
            color = CLR_TEXT
        elif state == "valid":
            add_rect(s, x, y + h - Inches(0.06), w, Pt(2), fill=CLR_BLUE, line=None)
            label = "●●●●●●●●●●"
        elif state == "invalid":
            add_rect(s, x, y + h - Inches(0.06), w, Pt(2), fill=CLR_RED, line=None)
            label = preset.get("text", "germweapon@gmail.com")
            # 에러 메시지
            add_text(s, x, y + h, w, Inches(0.18),
                     preset.get("error", "입력값을 확인해주세요."), size=8, color=CLR_RED)
        add_text(s, x, y + Inches(0.1), w, Inches(0.25), str(label), size=10, color=color)
    elif kind == "button":
        state = preset.get("state", "default")
        if state == "default":
            fill = CLR_GRAY_DARK
            txt_color = CLR_WHITE
        elif state == "disabled":
            fill = RGBColor(0xCB, 0xD5, 0xE1)
            txt_color = RGBColor(0x9C, 0xA3, 0xAF)
        else:
            fill = CLR_BLUE
            txt_color = CLR_WHITE
        add_rect(s, x, y + Inches(0.05), w, Inches(0.3), fill=fill, line=None, corner=True)
        add_text(s, x, y + Inches(0.08), w, Inches(0.3),
                 preset.get("label", "Button"), size=10, bold=True, color=txt_color, align=PP_ALIGN.CENTER)
    elif kind == "toggle":
        on = preset.get("state") == "on"
        track_w = Inches(0.6)
        track_h = Inches(0.22)
        track_x = x + (w - track_w) / 2
        track_y = y + Inches(0.1)
        add_rect(s, track_x, track_y, track_w, track_h,
                 fill=(CLR_BLUE if on else RGBColor(0xCB, 0xD5, 0xE1)), line=None, corner=True)
        knob_x = track_x + (track_w - Inches(0.18)) if on else track_x + Inches(0.02)
        add_oval(s, knob_x, track_y + Inches(0.02), Inches(0.18), Inches(0.18), fill=CLR_WHITE, line=None)
    else:
        add_rect(s, x, y + Inches(0.05), w, h - Inches(0.1), fill=CLR_BG_LIGHT, line=CLR_BORDER, corner=True)
        add_text(s, x, y + Inches(0.15), w, Inches(0.25),
                 preset.get("label", str(kind)), size=10, color=CLR_TEXT_MUTED, align=PP_ALIGN.CENTER)


def render_desc_table(s, x, y, w, title, items: list[dict]) -> Any:
    """Function / Button Description 표. 반환: 표 끝의 y 좌표."""
    header_h = Inches(0.3)
    add_rect(s, x, y, w, header_h, fill=CLR_HEADER_GRAY, line=None)
    add_text(s, x + Inches(0.15), y + Inches(0.05), w - Inches(0.3), header_h,
             title, size=11, bold=True, color=CLR_WHITE)
    cur = y + header_h
    marker_w = Inches(0.4)
    for item in items:
        marker = str(item.get("marker", item.get("id", "?")))
        title_text = item.get("title", "")
        lines = item.get("lines", [])
        # 행 높이 추정
        line_h = Inches(0.24)
        title_h = Inches(0.26)
        rows_h = title_h + line_h * len(lines) + Inches(0.1)
        # 외곽
        add_rect(s, x, cur, w, rows_h, fill=CLR_WHITE, line=CLR_BORDER)
        # 마커 셀
        add_rect(s, x, cur, marker_w, rows_h, fill=CLR_BG_LIGHT, line=CLR_BORDER)
        add_text(s, x, cur + Inches(0.05), marker_w, Inches(0.3),
                 marker, size=11, bold=True, align=PP_ALIGN.CENTER)
        # 본문
        body_x = x + marker_w + Inches(0.1)
        body_w = w - marker_w - Inches(0.2)
        add_text(s, body_x, cur + Inches(0.04), body_w, title_h, "• " + title_text, size=10, bold=True)
        cy = cur + title_h + Inches(0.02)
        for ln in lines:
            add_text(s, body_x + Inches(0.15), cy, body_w - Inches(0.15), line_h,
                     "- " + str(ln), size=9, color=CLR_TEXT)
            cy += line_h
        cur += rows_h
    return cur


def render_validation_table(s, x, y, w, rows: list[dict]) -> None:
    """우하 Validation 메시지 표."""
    cols_w = [Inches(0.5), Inches(1.0), Inches(2.2), Inches(0.9)]
    # 합이 4.6in과 정확히 안 맞으면 마지막 열로 보정
    total = sum(c for c in cols_w)
    if total < w:
        cols_w[-1] = cols_w[-1] + (w - total)
    headers = ["체크\n시점", "조건", "메시지", "노출유형"]
    header_h = Inches(0.4)
    # 헤더
    cx = x
    for i, h in enumerate(headers):
        add_rect(s, cx, y, cols_w[i], header_h, fill=CLR_TABLE_HEADER, line=CLR_BORDER)
        add_text(s, cx + Inches(0.05), y + Inches(0.05), cols_w[i] - Inches(0.1), header_h,
                 h, size=9, bold=True, color=CLR_WHITE, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    # 본문
    cur = y + header_h
    for row in rows:
        # 메시지 셀 내용 길이로 행 높이 추정
        message = row.get("message", "")
        title = row.get("message_title", "")
        buttons = row.get("message_buttons") or ([row["message_button"]] if row.get("message_button") else [])
        approx_chars = len(message) + (len(title) + 4 if title else 0) + sum(len(b) + 4 for b in buttons)
        n_lines = max(2, (approx_chars // 24) + 1)
        row_h = Inches(0.22 * n_lines + 0.2)

        cx = x
        # 체크시점
        add_rect(s, cx, cur, cols_w[0], row_h, fill=CLR_WHITE, line=CLR_BORDER)
        add_text(s, cx, cur + Inches(0.05), cols_w[0], Inches(0.3),
                 str(row.get("checkpoint", "")), size=9, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[0]
        # 조건
        add_rect(s, cx, cur, cols_w[1], row_h, fill=CLR_WHITE, line=CLR_BORDER)
        add_text(s, cx + Inches(0.05), cur + Inches(0.05), cols_w[1] - Inches(0.1), row_h,
                 str(row.get("condition", "")), size=9)
        cx += cols_w[1]
        # 메시지 — 팝업이면 ▶ + 제목 + 본문 + [버튼]
        add_rect(s, cx, cur, cols_w[2], row_h, fill=CLR_WHITE, line=CLR_BORDER)
        msg_lines: list[tuple[str, dict]] = []
        if title:
            msg_lines.append(("▶ " + title, {"size": 9, "bold": True}))
        msg_lines.append((message, {"size": 9}))
        if buttons:
            msg_lines.append(("[" + "] [".join(buttons) + "]", {"size": 9, "color": CLR_BLUE}))
        add_multiline(s, cx + Inches(0.05), cur + Inches(0.05), cols_w[2] - Inches(0.1), row_h, msg_lines)
        cx += cols_w[2]
        # 노출유형
        add_rect(s, cx, cur, cols_w[3], row_h, fill=CLR_WHITE, line=CLR_BORDER)
        add_text(s, cx + Inches(0.05), cur + Inches(0.05), cols_w[3] - Inches(0.1), row_h,
                 str(row.get("display", "")), size=9)
        cur += row_h


# ---------------------------------------------------------------------------
# 메인
# ---------------------------------------------------------------------------
def validate(data: dict) -> list[str]:
    errors = []
    meta = data.get("meta", {})
    if not meta.get("title"):
        errors.append("meta.title 누락")
    if not meta.get("version"):
        errors.append("meta.version 누락")
    ia = data.get("ia", {})
    boxes = {b["id"] for b in ia.get("boxes", [])}
    for c in ia.get("connections", []):
        if c.get("from") not in boxes:
            errors.append(f"connection from={c.get('from')}: 박스 없음")
        if c.get("to") not in boxes:
            errors.append(f"connection to={c.get('to')}: 박스 없음")
    return errors


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True)
    ap.add_argument("--output", required=True)
    ap.add_argument("--ia-render", choices=["shapes", "image"], default="shapes")
    ap.add_argument("--ia-image", default=None, help="--ia-render=image일 때 PNG 경로")
    ap.add_argument("--presets", default=None, help="component_state_presets.json 경로 (기본: 스크립트 옆 assets/)")
    args = ap.parse_args()

    with open(args.input, "r", encoding="utf-8") as f:
        data = json.load(f)

    doc_root = os.path.dirname(os.path.abspath(args.input))

    presets_path = args.presets
    if not presets_path:
        # 스크립트의 ../assets/component_state_presets.json
        here = os.path.dirname(os.path.abspath(__file__))
        candidate = os.path.join(here, "..", "assets", "component_state_presets.json")
        if os.path.exists(candidate):
            presets_path = candidate

    presets = {}
    if presets_path and os.path.exists(presets_path):
        with open(presets_path, "r", encoding="utf-8") as f:
            presets = json.load(f)

    errs = validate(data)
    if errs:
        print("VALIDATION ERRORS:", file=sys.stderr)
        for e in errs:
            print("  -", e, file=sys.stderr)
        # 진행은 함 (경고만)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs, data.get("meta", {}))
    slide_revision_history(prs, data.get("revision_history", []))
    slide_ia(prs, data, mode=args.ia_render, ia_image_path=args.ia_image)

    for screen in data.get("screens", []):
        slide_screen_detail(prs, screen, doc_root, presets)

    os.makedirs(os.path.dirname(os.path.abspath(args.output)) or ".", exist_ok=True)
    prs.save(args.output)
    print(f"Wrote PPTX: {args.output}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
